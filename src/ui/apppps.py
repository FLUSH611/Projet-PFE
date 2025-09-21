import os
import io
from pathlib import Path
import streamlit as st
import pandas as pd
from typing import List, Tuple

# ---- Optional libs (installed in your Step 0) ----
import chromadb
from chromadb.config import Settings
from sentence_transformers import SentenceTransformer
# We only use a light "LLM" summarizer if available; otherwise fallback to extractive display.
try:
    from transformers import pipeline, AutoTokenizer, AutoModelForSeq2SeqLM
    HF_AVAILABLE = True
except Exception:
    HF_AVAILABLE = False

# ---- PPT/DOCX generation ----
from pptx import Presentation
from docx import Document

# ---- Constants (paths) ----
ROOT = Path(__file__).resolve().parents[2] if (Path(__file__).resolve().parents[2] / "data").exists() else Path.cwd()
DATA_RAW = ROOT / "data" / "raw"
DATA_PROCESSED = ROOT / "data" / "processed"
VEC_DIR = ROOT / "vectors"
OUT_DIR = ROOT / "out"
COLLECTION_NAME = "itstorm_docs"
EMB_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"

# Ensure folders exist
for p in [DATA_RAW, DATA_PROCESSED, VEC_DIR, OUT_DIR]:
    p.mkdir(parents=True, exist_ok=True)

# ==========
#   STYLE
# ==========
st.set_page_config(
    page_title="Intelligent Consulting Copilot",
    page_icon="🧠",
    layout="wide"
)

CUSTOM_CSS = """
<style>
/* page */
section.main > div { padding-top: 1.2rem; }
.block-container{padding-top:1rem; padding-bottom:2rem;}
/* title + subtitle */
h1, h2, h3 { font-family: "Inter", system-ui, -apple-system, Segoe UI, Roboto, Arial; }
.big-title { font-size: 2.0rem; font-weight: 800; color: #0b285a; margin-bottom: 0.25rem; }
.sub-title { font-size: 0.95rem; color: #2b4162; opacity: 0.9; }
/* cards */
.kpi-card {
  border: 1px solid #e8edf7; border-radius: 16px;
  padding: 14px 16px; background: #fbfdff; box-shadow: 0 0 1px rgba(0,0,0,0.06);
}
.kpi-label { color:#5c6f91; font-size: 0.8rem; margin-bottom: 0.15rem; }
.kpi-value { color:#0b285a; font-size: 1.35rem; font-weight: 700; }
.card {
  border: 1px solid #e9eef7; border-radius: 16px; padding: 18px;
  background: white; box-shadow: 0 1px 2px rgba(10, 30, 80, 0.05);
}
/* buttons */
.stButton>button {
  border-radius: 10px; padding: 0.55rem 0.9rem; font-weight: 600;
  border: 1px solid #d9e3f0; background: #0b285a; color: #fff;
}
.stButton>button:hover { filter: brightness(1.06); }
/* chips */
.chip { display:inline-block; padding: 3px 9px; background:#eef3ff; border-radius: 9999px; color:#0b285a; font-size: 12px; margin-right:6px; }
.source { font-size: 12px; color:#4b628a; background: #f4f7fb; padding: 4px 8px; border-radius: 9999px; display:inline-block; margin-right:8px; margin-top:6px; }
hr { border: none; border-top: 1px solid #ecf0f7; margin: 0.8rem 0 1.0rem; }
</style>
"""
st.markdown(CUSTOM_CSS, unsafe_allow_html=True)

# ==========
#  HELPERS
# ==========
@st.cache_resource
def get_db():
    client = chromadb.Client(Settings(persist_directory=str(VEC_DIR), is_persistent=True))
    # ensure collection exists or create empty
    names = [c.name for c in client.list_collections()]
    coll = client.get_collection(COLLECTION_NAME) if COLLECTION_NAME in names else client.create_collection(COLLECTION_NAME)
    return client, coll

@st.cache_resource
def get_embedder():
    return SentenceTransformer(EMB_MODEL_NAME)

@st.cache_resource
def get_summarizer():
    if not HF_AVAILABLE:
        return None
    # Small, fast summarizer (pegasus_xsum or t5-small if unavailable)
    for name in ["google/pegasus-xsum", "t5-small"]:
        try:
            tok = AutoTokenizer.from_pretrained(name)
            mdl = AutoModelForSeq2SeqLM.from_pretrained(name)
            return pipeline("summarization", model=mdl, tokenizer=tok)
        except Exception:
            continue
    return None

def retrieve(query: str, k: int = 4, similarity_threshold: float = 0.0):
    _, coll = get_db()
    emb = get_embedder()
    q = emb.encode([query], normalize_embeddings=True).tolist()[0]
    res = coll.query(query_embeddings=[q], n_results=k)
    docs = res.get("documents", [[]])[0]
    metas = res.get("metadatas", [[]])[0]
    dists = res.get("distances", [[]])[0] if res.get("distances") else [0.0] * len(docs)

    # Filter if threshold is set and distances present (lower is closer in Chroma cosine? distances may be similarity metric)
    pairs = []
    for doc, meta, dist in zip(docs, metas, dists):
        if similarity_threshold > 0 and dist > similarity_threshold:
            continue
        pairs.append((doc, meta.get("source", "unknown"), dist))
    return pairs

def save_uploaded_files(files):
    saved = []
    for f in files:
        target = DATA_RAW / f.name
        with open(target, "wb") as out:
            out.write(f.read())
        saved.append(str(target))
    return saved

def extract_and_chunk(in_dir=DATA_RAW, out_path=DATA_PROCESSED / "chunks.jsonl", chunk_size=900, overlap=150):
    from pypdf import PdfReader
    import docx2txt, json

    def extract_pdf(p):
        reader = PdfReader(p)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    def extract_docx(p):
        return docx2txt.process(p)

    recs = []
    for p in Path(in_dir).glob("*"):
        text = ""
        if p.suffix.lower() == ".pdf":
            text = extract_pdf(p)
        elif p.suffix.lower() == ".docx":
            text = extract_docx(p)
        elif p.suffix.lower() in {".txt", ".md"}:
            text = p.read_text(encoding="utf-8", errors="ignore")
        if not text:
            continue
        words = text.split()
        i = 0
        while i < len(words):
            chunk_words = words[i:i+chunk_size]
            recs.append({"source": str(p), "text": " ".join(chunk_words)})
            i += (chunk_size - overlap)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    with open(out_path, "w", encoding="utf-8") as f:
        for r in recs:
            f.write(__import__("json").dumps(r, ensure_ascii=False) + "\n")
    return len(recs), out_path

def rebuild_index(in_path=DATA_PROCESSED / "chunks.jsonl"):
    client, _ = get_db()
    # reset collection
    try:
        client.delete_collection(COLLECTION_NAME)
    except Exception:
        pass
    coll = client.create_collection(COLLECTION_NAME)
    model = get_embedder()

    texts, metas, ids = [], [], []
    import json
    with open(in_path, "r", encoding="utf-8") as f:
        for k, line in enumerate(f):
            r = json.loads(line)
            texts.append(r["text"])
            metas.append({"source": r["source"]})
            ids.append(f"id_{k}")

    emb = model.encode(texts, normalize_embeddings=True).tolist()
    coll.add(documents=texts, embeddings=emb, metadatas=metas, ids=ids)
    return len(texts)

def kpi_card(label: str, value: str):
    st.markdown(f"""
    <div class="kpi-card">
      <div class="kpi-label">{label}</div>
      <div class="kpi-value">{value}</div>
    </div>
    """, unsafe_allow_html=True)

def render_sources(rows: List[Tuple[str, str, float]]):
    if not rows:
        st.info("No sources retrieved yet.")
        return
    for i, (doc, src, dist) in enumerate(rows, 1):
        with st.expander(f"Context #{i} — {Path(src).name}"):
            st.write(doc[:1200] + ("..." if len(doc) > 1200 else ""))
            st.markdown(f'<span class="source">Source: {src}</span>', unsafe_allow_html=True)

def summarize_text(text: str, max_chars=1200):
    if not text:
        return ""
    text = text[:4000]  # safety
    if get_summarizer() is None:
        # Fallback: simple heuristic summary
        parts = [p.strip() for p in text.split("\n") if p.strip()]
        return " ".join(parts[:4])[:max_chars]
    summarizer = get_summarizer()
    try:
        out = summarizer(text, max_length=128, min_length=50, do_sample=False)
        return out[0]["summary_text"]
    except Exception:
        return text[:max_chars]

def make_ppt(title="Client Update", bullets=None) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if bullets:
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Key Points"
        tf = s2.shapes.placeholders[1].text_frame
        tf.clear()
        for b in bullets:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.read()

def make_docx(title="Executive Summary", paragraphs=None) -> bytes:
    doc = Document()
    doc.add_heading(title, 0)
    if paragraphs:
        for p in paragraphs:
            doc.add_paragraph(p)
    bio = io.BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio.read()

# ==========
#  HEADER
# ==========
col1, col2 = st.columns([0.8, 0.2])
with col1:
    st.markdown('<div class="big-title">Intelligent Consulting Copilot</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-title">Knowledge Navigator · Auto-Deliverables · Market Intelligence</div>', unsafe_allow_html=True)
with col2:
    st.markdown('<div style="text-align:right;"><span class="chip">GenAI</span><span class="chip">RAG</span><span class="chip">Consulting</span></div>', unsafe_allow_html=True)

st.write("")
k1, k2, k3, k4 = st.columns(4)
with k1: kpi_card("Indexed Docs", f"{len(list(DATA_RAW.glob('*')))} files")
with k2: kpi_card("Chunks File", "✅" if (DATA_PROCESSED / "chunks.jsonl").exists() else "❌")
with k3: kpi_card("Vector DB", "✅" if any(VEC_DIR.glob('*')) else "❌")
with k4: kpi_card("Output Folder", "✅" if any(OUT_DIR.glob('*')) else "—")

# ==========
#  SIDEBAR
# ==========
st.sidebar.header("Settings")
top_k = st.sidebar.slider("Top-K results", 2, 12, 4, 1)
sim_thresh = st.sidebar.slider("Similarity threshold (distance)", 0.0, 1.5, 0.0, 0.05)
model_name = st.sidebar.text_input("Embedding model", EMB_MODEL_NAME)
if model_name != EMB_MODEL_NAME:
    st.sidebar.warning("This UI uses MiniLM; change in code to switch model safely.")

st.sidebar.markdown("---")
st.sidebar.caption("Theme: Navy / Consulting · v1.0")

# ==========
#  TABS
# ==========
tab1, tab2, tab3, tab4 = st.tabs(["💬 Knowledge Chat", "📂 Upload & Index", "📝 Generate Docs", "🌍 Market Watch"])

# ---- TAB 1: CHAT ----
with tab1:
    st.markdown("### Knowledge Q&A")
    q = st.text_input("Ask a question about your indexed documents")
    go = st.button("🔎 Retrieve & Answer", use_container_width=True)
    if go and q:
        with st.spinner("Retrieving relevant context..."):
            rows = retrieve(q, k=top_k, similarity_threshold=sim_thresh)
        if not rows:
            st.warning("No relevant context found. Try another query or re-index.")
        else:
            # Build synthetic answer using retrieved text (short summary)
            joined = "\n\n".join([r[0] for r in rows])
            answer = summarize_text(joined)
            st.markdown("#### Answer")
            st.write(answer if answer else "Context retrieved, but no summary available.")
            st.markdown("#### Sources")
            render_sources(rows)

# ---- TAB 2: UPLOAD & INDEX ----
with tab2:
    st.markdown("### Upload documents & (re)build index")
    files = st.file_uploader("Upload PDF/DOCX/TXT", type=["pdf", "docx", "txt", "md"], accept_multiple_files=True)
    c1, c2 = st.columns(2)
    if c1.button("⬆️ Save uploads", use_container_width=True) and files:
        saved = save_uploaded_files(files)
        st.success(f"Saved {len(saved)} file(s) to `data/raw/`.")
    if c2.button("⚙️ Extract & Chunk", use_container_width=True):
        with st.spinner("Extracting & chunking..."):
            n, outp = extract_and_chunk()
        st.success(f"Created {n} chunks → {outp}")
    if st.button("🧠 Rebuild Vector Index", use_container_width=True):
        with st.spinner("Encoding & indexing... this can take a moment"):
            total = rebuild_index()
        st.success(f"Indexed {total} chunks into ChromaDB.")

    # Quick preview
    if (DATA_PROCESSED / "chunks.jsonl").exists():
        st.markdown("#### Preview extracted chunks")
        # display only first 3 lines to avoid heavy rendering
        import json, itertools
        rows = []
        with open(DATA_PROCESSED / "chunks.jsonl", "r", encoding="utf-8") as f:
            for line in itertools.islice(f, 3):
                rows.append(json.loads(line))
        df = pd.DataFrame(rows)
        st.dataframe(df, use_container_width=True, height=200)

# ---- TAB 3: DOCS GENERATION ----
with tab3:
    st.markdown("### Generate consulting deliverables")
    default_title = "Executive Summary – Client X"
    title = st.text_input("Document Title", value=default_title)
    st.caption("Tip: Paste text below or run a query in the Chat tab and copy the answer here.")
    txt = st.text_area("Content (will be summarized if too long)", height=220)

    c1, c2 = st.columns(2)
    if c1.button("📄 Download DOCX", use_container_width=True):
        bullets = [b.strip() for b in txt.split("\n") if b.strip()]
        paragraphs = bullets if bullets else ["No content provided."]
        data = make_docx(title=title, paragraphs=paragraphs)
        st.download_button("⬇️ Save DOCX", data=data, file_name="report.docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
    if c2.button("📊 Download PPTX", use_container_width=True):
        bullets = [b.strip() for b in txt.split("\n") if b.strip()]
        data = make_ppt(title=title, bullets=bullets[:8] if bullets else ["No content provided."])
        st.download_button("⬇️ Save PPTX", data=data, file_name="deck.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)

# ---- TAB 4: MARKET WATCH (placeholder) ----
with tab4:
    st.markdown("### Competitive & Market Watch")
    st.info("This is a placeholder dashboard. In Month 3, wire your scraping (Scrapy/Playwright) to auto-populate.")
    st.markdown("**Suggested sections:**")
    st.markdown("- New tenders / RFPs detected")
    st.markdown("- Competitor news & hiring")
    st.markdown("- Tech trends (GenAI, Cloud, Data)")
    st.markdown("- Weekly PDF auto-report")

    # Minimal mock table for the feel
    import pandas as pd
    import datetime as dt
    import io
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    # --- KPIs row ---
    c1, c2, c3, c4 = st.columns(4)
    with c1: kpi_card("New items (7d)", "12")
    with c2: kpi_card("Opportunities", "4")
    with c3: kpi_card("Competitor Signals", "6")
    with c4: kpi_card("Tech Trends", "2")

    # --- Demo data (replace with your scraped feed later) ---
    demo = pd.DataFrame([
        {"date": (dt.date.today() - dt.timedelta(days=1)).isoformat(), "category": "RFP", "source": "gov.tenders", "title": "Cloud Migration RFP – Public Agency", "score": 0.86},
        {"date": (dt.date.today() - dt.timedelta(days=2)).isoformat(), "category": "Competitor", "source": "LinkedIn", "title": "Competitor X hiring 10 Data Engineers", "score": 0.74},
        {"date": (dt.date.today() - dt.timedelta(days=3)).isoformat(), "category": "Trend", "source": "TechCrunch", "title": "GenAI adoption in enterprise surges", "score": 0.69},
        {"date": (dt.date.today() - dt.timedelta(days=5)).isoformat(), "category": "RFP", "source": "eu.tenders", "title": "Data Platform Modernization – RFP", "score": 0.81},
        {"date": (dt.date.today() - dt.timedelta(days=6)).isoformat(), "category": "Competitor", "source": "Company PR", "title": "Competitor Y wins major AI project", "score": 0.77},
    ])

    # Filters
    st.markdown("#### Filters")
    c1, c2, c3 = st.columns([1,1,2])
    with c1:
        cat = st.multiselect("Category", options=sorted(demo["category"].unique()), default=list(demo["category"].unique()))
    with c2:
        min_score = st.slider("Min. relevance score", 0.0, 1.0, 0.65, 0.01)
    with c3:
        keyword = st.text_input("Keyword contains", value="")

    df = demo.copy()
    if cat:
        df = df[df["category"].isin(cat)]
    df = df[df["score"] >= min_score]
    if keyword.strip():
        df = df[df["title"].str.contains(keyword, case=False, na=False)]

    st.markdown("#### This week’s signals")
    st.dataframe(df.sort_values("date", ascending=False), use_container_width=True, height=260)

    # --- Upload CSV to simulate scraped feed ---
    st.markdown("#### Load external CSV (simulated feed)")
    csv_file = st.file_uploader("Upload a CSV with columns: date,category,source,title,score", type=["csv"])
    if csv_file is not None:
        try:
            ext_df = pd.read_csv(csv_file)
            st.success(f"Loaded {len(ext_df)} rows from CSV.")
            st.dataframe(ext_df.head(10), use_container_width=True, height=240)
        except Exception as e:
            st.error(f"Failed to read CSV: {e}")

    # --- Quick frequency chart (by category) ---
    st.markdown("#### Category distribution (this week)")
    freq = df["category"].value_counts().reset_index()
    freq.columns = ["category", "count"]
    if not freq.empty:
        st.bar_chart(freq.set_index("category"))

    # --- Export a Weekly PDF summary (dummy) ---
    st.markdown("#### Weekly PDF auto-report")
    if st.button("📄 Generate Weekly Report (PDF)", use_container_width=True):
        # Build a very small PDF summary for demo
        buffer = io.BytesIO()
        c = canvas.Canvas(buffer, pagesize=A4)
        width, height = A4
        y = height - 60
        c.setFont("Helvetica-Bold", 14)
        c.drawString(40, y, "Weekly Market Watch Report")
        y -= 24
        c.setFont("Helvetica", 10)
        c.drawString(40, y, f"Generated: {dt.datetime.now().strftime('%Y-%m-%d %H:%M')}")
        y -= 30

        # KPIs
        kpis = [
            ("New items (7d)", str(len(demo))),
            ("Opportunities", str((demo["category"] == "RFP").sum())),
            ("Competitor Signals", str((demo["category"] == "Competitor").sum())),
            ("Tech Trends", str((demo["category"] == "Trend").sum())),
        ]
        for label, val in kpis:
            c.drawString(40, y, f"- {label}: {val}")
            y -= 16

        y -= 12
        c.setFont("Helvetica-Bold", 12)
        c.drawString(40, y, "Top signals")
        y -= 20
        c.setFont("Helvetica", 10)

        for _, row in df.sort_values(["score", "date"], ascending=[False, False]).head(8).iterrows():
            line = f"[{row['date']}] ({row['category']}) {row['title']}  — score {row['score']:.2f}"
            c.drawString(40, y, line[:110])
            y -= 14
            if y < 80:
                c.showPage()
                y = height - 60
                c.setFont("Helvetica", 10)

        c.showPage()
        c.save()
        buffer.seek(0)
        st.download_button(
            "⬇️ Download weekly_report.pdf",
            data=buffer.getvalue(),
            file_name="weekly_report.pdf",
            mime="application/pdf",
            use_container_width=True
        )

    st.caption("Tip: Replace this placeholder with your Scrapy/Playwright pipeline in Month 3 and write into CSV/DB, then load & visualize here.")

